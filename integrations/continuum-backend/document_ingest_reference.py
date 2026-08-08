"""
Reference: extend Render /memories/ingest to extract text from Office documents.

Deploy on continuum-backend (FastAPI). Requires:
  pip install python-docx python-pptx openpyxl

Wire into existing ingest handler before chunking/embedding.
"""

from __future__ import annotations

import io
import os
from typing import Optional


def extract_document_text(filename: str, raw: bytes) -> Optional[str]:
    name = (filename or "").lower()
    if name.endswith(".docx"):
        from docx import Document

        doc = Document(io.BytesIO(raw))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    if name.endswith(".doc"):
        # Legacy binary .doc (OLE2). The app now extracts text on-device for chat,
        # but Layer 5 ingest should also handle .doc so users can upload old files.
        # Install one of the following on the Render host:
        #   pip install textract        (wraps antiword; needs antiword binary)
        #   apt-get install antiword
        #   OR use LibreOffice headless: soffice --headless --convert-to txt --outdir /tmp file.doc
        try:
            import textract  # type: ignore

            with open("/tmp/_continuum_doc.doc", "wb") as fh:
                fh.write(raw)
            text = textract.process("/tmp/_continuum_doc.doc").decode("utf-8", errors="replace")
            os.remove("/tmp/_continuum_doc.doc")
            return text.strip() or None
        except Exception:
            return None
    if name.endswith(".pptx"):
        from pptx import Presentation

        prs = Presentation(io.BytesIO(raw))
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text.strip())
        return "\n".join(lines)
    if name.endswith(".ppt"):
        return None
    if name.endswith(".xlsx"):
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        lines = []
        for sheet in wb.worksheets:
            lines.append(f"## Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if cells:
                    lines.append("\t".join(cells))
        wb.close()
        return "\n".join(lines)
    if name.endswith(".xls"):
        return None
    if name.endswith(".txt"):
        return raw.decode("utf-8", errors="replace")
    return None
